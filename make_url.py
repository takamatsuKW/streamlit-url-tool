import streamlit as st
import pyshorteners
import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import logging

logging.getLogger("streamlit").setLevel(logging.ERROR)

# 短縮URLを作成する関数
def create_short_url(long_url):
    s = pyshorteners.Shortener()
    return s.tinyurl.short(long_url)

# URLをQRコード（画像）にする関数
def create_qr_code(url):
    qr = qrcode.QRCode(box_size=10, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    return img

# PowerPointを作成する関数
def create_ppt(short_url, qr_img):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # タイトル追加
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "アンケートのお願い"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # QRコード画像をメモリに保存
    qr_stream = BytesIO()
    qr_img.save(qr_stream, format='PNG')
    qr_stream.seek(0)

    # QRコード画像追加
    slide.shapes.add_picture(qr_stream, Inches(3), Inches(2), Inches(4), Inches(4))

    # 短縮URLを追加
    text_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = short_url
    text_frame.paragraphs[0].font.size = Pt(20)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # PowerPointをメモリに保存
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    return ppt_stream

# Streamlitアプリ
st.title("URL短縮 & QRコード生成")

# ユーザー入力
long_url = st.text_input("短縮したいURLを入力してください:")

if long_url:
    # 短縮URL作成
    short_url = create_short_url(long_url)
    st.markdown(f"### 短縮URL: [ {short_url} ]({short_url})")

    # QRコード作成
    qr_img = create_qr_code(short_url)
    qr_stream = BytesIO()
    qr_img.save(qr_stream, format='PNG')
    qr_stream.seek(0)

    # QRコード表示
    st.image(qr_stream, caption="QRコード", use_column_width=False)

    # QRコードダウンロード
    st.download_button(
        label="QRコードをダウンロード",
        data=qr_stream,
        file_name="qr_code.png",
        mime="image/png"
    )

    # PowerPoint作成
    ppt_stream = create_ppt(short_url, qr_img)

    # PowerPointダウンロード
    st.download_button(
        label="PPTXをダウンロード",
        data=ppt_stream,
        file_name="アンケート_スライド.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
